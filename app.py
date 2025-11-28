import os
import streamlit as st
import pandas as pd
from PIL import Image
import docx
import pptx
import tempfile
import time
import concurrent.futures # For Parallel Processing

# Langchain & AI Imports
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_google_genai import GoogleGenerativeAIEmbeddings
import google.generativeai as genai
from langchain_community.vectorstores import FAISS
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
from langchain_experimental.agents import create_pandas_dataframe_agent
from dotenv import load_dotenv

# Load Environment
load_dotenv()
os.environ["GOOGLE_API_KEY"] = os.getenv("GOOGLE_API_KEY")
genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))

# ------------------------------------------------------------------------
# 0. UI STYLING (3D SPACE & ANIMATIONS)
# ------------------------------------------------------------------------
def set_custom_ui():
    st.markdown("""
    <style>
        /* 1. IMPORT FONTS */
        @import url('https://fonts.googleapis.com/css2?family=Orbitron:wght@400;700&family=Roboto:wght@300;400&display=swap');

        /* 2. THE 3D SPACE BACKGROUND */
        .stApp {
            background: radial-gradient(ellipse at bottom, #1b2735 0%, #090a0f 100%);
            color: #fff;
            font-family: 'Roboto', sans-serif;
        }

        /* Creating stars using CSS */
        #stars {
            width: 1px; height: 1px;
            background: transparent;
            box-shadow: 1744px 122px #FFF , 134px 1321px #FFF , 92px 859px #FFF; /* This continues for many pixels */
            animation: animStar 50s linear infinite;
        }
        
        /* 3. 3D GLOWING BUTTONS */
        .stButton>button {
            background: linear-gradient(145deg, #1e232a, #2a303a);
            color: #4b90ff;
            border: 2px solid #4b90ff;
            border-radius: 12px;
            box-shadow: 
                5px 5px 10px #13171d, 
                -5px -5px 10px #353d49,
                0 0 10px #4b90ff; /* Glow */
            transition: all 0.3s ease;
            font-family: 'Orbitron', sans-serif;
            text-transform: uppercase;
            letter-spacing: 2px;
            font-weight: bold;
        }

        .stButton>button:hover {
            transform: translateY(-5px) scale(1.02);
            box-shadow: 
                0 10px 20px rgba(75, 144, 255, 0.4),
                0 0 20px #4b90ff inset;
            color: #fff;
            border-color: #fff;
        }

        .stButton>button:active {
            transform: translateY(2px);
            box-shadow: inset 5px 5px 10px #13171d, inset -5px -5px 10px #353d49;
        }

        /* 4. CHAT MESSAGE STYLING (GLASSMORPHISM) */
        .stChatMessage {
            background: rgba( 255, 255, 255, 0.05 );
            box-shadow: 0 8px 32px 0 rgba( 31, 38, 135, 0.37 );
            backdrop-filter: blur( 4px );
            -webkit-backdrop-filter: blur( 4px );
            border-radius: 10px;
            border: 1px solid rgba( 255, 255, 255, 0.18 );
            margin-bottom: 15px;
        }
        
        /* User Message specific */
        div[data-testid="stChatMessage"]:nth-child(odd) {
            background: rgba( 75, 144, 255, 0.1 );
            border-left: 5px solid #4b90ff;
        }

        /* Sidebar Styling */
        section[data-testid="stSidebar"] {
            background: #0d1117;
            border-right: 1px solid #333;
        }
        
        /* Headers */
        h1, h2, h3 {
            font-family: 'Orbitron', sans-serif;
            text-shadow: 0 0 10px #4b90ff;
        }

        /* SPINNER ANIMATION (Thinking...) */
        .stSpinner > div {
            border-top-color: #4b90ff !important;
        }

    </style>
    """, unsafe_allow_html=True)

# ------------------------------------------------------------------------
# 1. FAST PROCESSING FUNCTIONS (PARALLEL)
# ------------------------------------------------------------------------

def get_gemini_vision_text(file_data):
    """Worker function for parallel processing."""
    file_bytes, mime_type, file_name = file_data
    with tempfile.NamedTemporaryFile(delete=False, suffix=f".{mime_type.split('/')[-1]}") as temp_file:
        temp_file.write(file_bytes)
        temp_file_path = temp_file.name

    try:
        myfile = genai.upload_file(temp_file_path, mime_type=mime_type)
        while myfile.state.name == "PROCESSING":
            time.sleep(1)
            myfile = genai.get_file(myfile.name)

        model = genai.GenerativeModel("gemini-2.5-flash")
        prompt = "Transcribe ALL text, handwriting, and describe charts/graphs in detail. Return plain text."
        response = model.generate_content([myfile, prompt])
        
        myfile.delete()
        os.remove(temp_file_path)
        return f"\n--- File: {file_name} ---\n{response.text}\n"

    except Exception as e:
        return f"\n[Error processing {file_name}: {str(e)}]\n"

def process_uploaded_files_parallel(uploaded_files):
    raw_text = ""
    excel_data = None
    vision_tasks = [] 
    
    for file in uploaded_files:
        file_bytes = file.read()
        file_name = file.name
        
        if file_name.endswith('.xlsx'):
            try:
                file.seek(0)
                df = pd.read_excel(file)
                excel_data = df
                raw_text += f"\n[Excel Data: {file_name}]\n{df.to_markdown(index=False)}\n"
            except Exception: pass
        
        elif file_name.endswith('.docx'):
            file.seek(0)
            doc = docx.Document(file)
            text = "\n".join([p.text for p in doc.paragraphs])
            raw_text += f"\n[Word Doc: {file_name}]\n{text}\n"

        elif file_name.endswith('.pptx'):
            file.seek(0)
            prs = pptx.Presentation(file)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            raw_text += f"\n[PowerPoint: {file_name}]\n{text}\n"

        elif file_name.endswith(('.pdf', '.png', '.jpg', '.jpeg')):
            mime_type = "application/pdf" if file_name.endswith('.pdf') else file.type
            vision_tasks.append((file_bytes, mime_type, file_name))

    if vision_tasks:
        with concurrent.futures.ThreadPoolExecutor(max_workers=5) as executor:
            results = list(executor.map(get_gemini_vision_text, vision_tasks))
            for res in results:
                raw_text += res

    return raw_text, excel_data

# ------------------------------------------------------------------------
# 2. RAG & AGENT LOGIC
# ------------------------------------------------------------------------

def get_text_chunks(text):
    splitter = RecursiveCharacterTextSplitter(chunk_size=10000, chunk_overlap=1000)
    chunks = splitter.split_text(text)
    return chunks

def get_vector_store(chunks):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    vector_store = FAISS.from_texts(chunks, embedding=embeddings)
    vector_store.save_local("faiss_index")

def rag_query(user_question):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    new_db = FAISS.load_local("faiss_index", embeddings, allow_dangerous_deserialization=True)
    docs = new_db.similarity_search(user_question)
    
    model = ChatGoogleGenerativeAI(model="gemini-2.5-flash", temperature=0.3)
    prompt_template = """
    Answer the question in a helpful, friendly way. Use the context provided.
    
    Context:\n {context}?\n
    Question: \n{question}\n

    Answer:
    """
    prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])
    chain = load_qa_chain(llm=model, chain_type="stuff", prompt=prompt)
    response = chain({"input_documents": docs, "question": user_question}, return_only_outputs=True)
    return response['output_text']

def query_pandas_agent(df, question):
    llm = ChatGoogleGenerativeAI(model="gemini-2.5-flash", temperature=0)
    agent = create_pandas_dataframe_agent(llm, df, verbose=True, allow_dangerous_code=True)
    instructions = "Use matplotlib/seaborn for plots. Save as 'temp_plot.png'. Do not show."
    try:
        return agent.run(instructions + question)
    except Exception as e:
        return f"Analysis Error: {e}"

# ------------------------------------------------------------------------
# 3. MAIN UI APP
# ------------------------------------------------------------------------

def main():
    st.set_page_config(page_title="DOC-CHAT AI", page_icon="🌌", layout="wide")
    set_custom_ui() # APPLY 3D STYLES

    # Sidebar
    with st.sidebar:
        st.markdown("### 🌌 3D Document Hub")
        st.info("Upload PDFs, Handwriting, Excel, Word or Images.")
        
        uploaded_files = st.file_uploader("Upload Files", accept_multiple_files=True, label_visibility="collapsed")
        
        if st.button("🚀 IGNITE PROCESSING"):
            if not uploaded_files:
                st.warning("Please upload files first.")
            else:
                # Custom 3D Spinner Text
                with st.spinner("✨ Initializing Hyper-Drive Parallel Processing..."):
                    raw_text, excel_data = process_uploaded_files_parallel(uploaded_files)
                    
                    if excel_data is not None:
                        st.session_state.excel_data = excel_data
                    
                    if raw_text:
                        chunks = get_text_chunks(raw_text)
                        get_vector_store(chunks)
                        st.success("✅ Neural Network Updated!")
                    else:
                        st.warning("No readable text found.")
        
        st.markdown("---")
        if st.button("🗑️ PURGE MEMORY"):
            st.session_state.messages = []

    # Main Area
    st.markdown('<h1 style="font-size: 3em;">🤖 DOC-CHAT  <span style="color:#4b90ff;">AI</span></h1>', unsafe_allow_html=True)
    st.markdown("#### Ready to analyze your data in hyperspace.")

    # Initialize Chat
    if "messages" not in st.session_state:
        st.session_state.messages = [{"role": "assistant", "content": "System Online. Upload documents to begin analysis."}]
    if "excel_data" not in st.session_state:
        st.session_state.excel_data = None

    # Display History
    for message in st.session_state.messages:
        with st.chat_message(message["role"], avatar="🤖" if message["role"] == "assistant" else "👨‍🚀"):
            st.write(message["content"])

    # Chat Input
    if prompt := st.chat_input("Input command for analysis..."):
        st.session_state.messages.append({"role": "user", "content": prompt})
        with st.chat_message("user", avatar="👨‍🚀"):
            st.write(prompt)

        with st.chat_message("assistant", avatar="🤖"):
            message_placeholder = st.empty()
            
            # 3D Thinking UI
            with st.spinner("Processing in Quantum Space..."):
                # Logic
                is_plot = any(x in prompt.lower() for x in ["plot", "graph", "chart"]) and st.session_state.excel_data is not None
                
                if is_plot:
                    response = query_pandas_agent(st.session_state.excel_data, prompt)
                    if os.path.exists("temp_plot.png"):
                        st.image("temp_plot.png", caption="Visual Data Output")
                        os.remove("temp_plot.png")
                else:
                    response = rag_query(prompt)
                
                message_placeholder.markdown(response)
                
        st.session_state.messages.append({"role": "assistant", "content": response})

if __name__ == "__main__":
    main()